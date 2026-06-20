"""Generate DealRoom AI presentation with full brand styling."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Brand palette ─────────────────────────────────────────────────────────────
CREAM    = RGBColor(0xFA, 0xF9, 0xF5)
GREEN    = RGBColor(0x1A, 0x5E, 0x3A)
GOLD     = RGBColor(0xD4, 0xA8, 0x4B)
INK      = RGBColor(0x14, 0x21, 0x1A)
SAND     = RGBColor(0xE8, 0xE2, 0xD0)
STONE    = RGBColor(0xB8, 0xAF, 0x9C)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xF4, 0xF2, 0xED)
DKGREEN  = RGBColor(0x0D, 0x3D, 0x26)

# Slide size: 16:9
W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ── Low-level helpers ─────────────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=Pt(18), bold=False, color=INK,
             align=PP_ALIGN.LEFT, italic=False, font="Helvetica Neue"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name   = font
    return txb


def add_multiline(slide, lines, x, y, w, h, size=Pt(16), color=INK,
                  spacing=Pt(6), font="Helvetica Neue"):
    """lines = list of (text, bold, size_override_or_None)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for (txt, bold, sz) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        run = p.add_run()
        run.text = txt
        run.font.size  = sz or size
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.name  = font
    return txb


def add_bullet(slide, items, x, y, w, h, size=Pt(17), color=INK,
               bullet_color=GOLD, font="Helvetica Neue", spacing=Pt(8)):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        # bullet char
        run0 = p.add_run()
        run0.text = "▸  "
        run0.font.size  = size
        run0.font.color.rgb = bullet_color
        run0.font.name  = font
        run = p.add_run()
        run.text = item
        run.font.size  = size
        run.font.color.rgb = color
        run.font.name  = font
    return txb


# ── Logo drawn as shapes ──────────────────────────────────────────────────────

def draw_logo(slide, cx, cy, scale=1.0):
    """Draw the split-page logo icon at center (cx,cy) with optional scale."""
    u = Inches(0.04) * scale  # 1 unit

    icon_w = u * 40
    icon_h = u * 80
    ix = cx - icon_w / 2
    iy = cy - icon_h / 2

    hw = icon_w / 2

    # Left half (sand)
    add_rect(slide, ix, iy, hw, icon_h, fill=SAND)
    # Right half (green)
    add_rect(slide, ix + hw, iy, hw, icon_h, fill=GREEN)

    # Lines on left
    lx1, lx2 = ix + u*4, ix + hw - u*4
    for frac in [0.2, 0.33, 0.46, 0.59, 0.72]:
        add_rect(slide, lx1, iy + icon_h*frac, lx2-lx1, u*2.2, fill=STONE)

    # Lines on right (shorter, white)
    rx1 = ix + hw + u*4
    rx2 = ix + icon_w - u*4
    for i, frac in enumerate([0.2, 0.33, 0.46]):
        w_ = (rx2-rx1) * (1 - i*0.12)
        add_rect(slide, rx1, iy + icon_h*frac, w_, u*2.2, fill=CREAM)

    # Gold circle checkmark area (bottom right quadrant)
    cr = u*9
    ccx = ix + hw + hw/2
    ccy = iy + icon_h*0.77

    # Draw circle as ellipse
    ell = slide.shapes.add_shape(9, ccx-cr, ccy-cr, cr*2, cr*2)  # 9 = oval
    ell.fill.background()
    ell.line.color.rgb = GOLD
    ell.line.width = int(Pt(2.5) * scale)

    # Checkmark via two rectangles (simplified)
    # left arm: small rect angled - draw as thin rect
    ck_x1 = ccx - cr*0.5
    ck_y1 = ccy + cr*0.1
    add_rect(slide, ck_x1, ck_y1 - u, u*2.5, u*2.5, fill=GOLD)
    add_rect(slide, ccx - cr*0.5, ccy + cr*0.05, cr*1.2, u*2.5, fill=GOLD)


# ── Section header bar ────────────────────────────────────────────────────────

def section_bar(slide, label):
    """Green bar with label at the bottom."""
    add_rect(slide, Inches(0), H - Inches(0.42), W, Inches(0.42), fill=GREEN)
    add_text(slide, label, Inches(0.3), H - Inches(0.42), W - Inches(0.6),
             Inches(0.42), size=Pt(11), color=CREAM, align=PP_ALIGN.LEFT)


def gold_line(slide, x, y, w, thick=Pt(1.2)):
    add_rect(slide, x, y, w, thick, fill=GOLD)


def slide_title(slide, title, subtitle=None, x=Inches(0.7), y=Inches(0.4),
                w=Inches(11), title_size=Pt(34)):
    add_text(slide, title, x, y, w, Inches(0.75),
             size=title_size, bold=True, color=GREEN, font="Georgia")
    if subtitle:
        gold_line(slide, x, y + Inches(0.75), w * 0.35)
        add_text(slide, subtitle, x, y + Inches(0.92), w, Inches(0.45),
                 size=Pt(15), color=STONE, font="Helvetica Neue")


# ═════════════════════════════════════════════════════════════════════════════
#   BUILD SLIDES
# ═════════════════════════════════════════════════════════════════════════════

prs = new_prs()
blank = prs.slide_layouts[6]  # completely blank

# ── 1. TITLE SLIDE ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)

# Dark green left band
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)
add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=GOLD)

# Logo icon (centered-left)
draw_logo(s, Inches(3.4), Inches(3.3), scale=1.6)

# Wordmark
add_text(s, "DealRoom", Inches(4.4), Inches(2.1), Inches(5.5), Inches(0.9),
         size=Pt(58), bold=True, color=INK, font="Georgia")
add_text(s, "AI", Inches(7.95), Inches(2.1), Inches(2), Inches(0.9),
         size=Pt(58), bold=False, color=GREEN, font="Georgia")

# Tagline
gold_line(s, Inches(4.4), Inches(3.05), Inches(5))
add_text(s, "DUE DILIGENCE.  REIMAGINED.", Inches(4.4), Inches(3.2), Inches(8), Inches(0.45),
         size=Pt(13), bold=True, color=INK, font="Helvetica Neue")

# Subtitle
add_text(s, "AI-Powered M&A Due Diligence Platform",
         Inches(4.4), Inches(3.78), Inches(8), Inches(0.5),
         size=Pt(19), bold=False, color=GREEN, font="Georgia")

# Presenter / date placeholder
add_text(s, "2026", Inches(4.4), Inches(6.7), Inches(8), Inches(0.4),
         size=Pt(12), color=STONE)

section_bar(s, "DealRoom AI  ·  Confidential")


# ── 2. THE PROBLEM ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "The Problem", "Traditional M&A due diligence is broken")

# Big stat box
add_rect(s, Inches(0.7), Inches(1.55), Inches(3.5), Inches(1.4), fill=GREEN)
add_text(s, "6–12 Weeks", Inches(0.85), Inches(1.65), Inches(3.2), Inches(0.7),
         size=Pt(32), bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="Georgia")
add_text(s, "average time for a full DD report",
         Inches(0.85), Inches(2.25), Inches(3.2), Inches(0.5),
         size=Pt(13), color=CREAM, align=PP_ALIGN.CENTER)

add_rect(s, Inches(4.5), Inches(1.55), Inches(3.5), Inches(1.4), fill=GREEN)
add_text(s, "10+ Analysts", Inches(4.65), Inches(1.65), Inches(3.2), Inches(0.7),
         size=Pt(32), bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="Georgia")
add_text(s, "reading thousands of pages manually",
         Inches(4.65), Inches(2.25), Inches(3.2), Inches(0.5),
         size=Pt(13), color=CREAM, align=PP_ALIGN.CENTER)

add_rect(s, Inches(8.3), Inches(1.55), Inches(3.5), Inches(1.4), fill=GREEN)
add_text(s, "~30% of deals", Inches(8.45), Inches(1.65), Inches(3.2), Inches(0.7),
         size=Pt(32), bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="Georgia")
add_text(s, "miss critical risks during diligence",
         Inches(8.45), Inches(2.25), Inches(3.2), Inches(0.5),
         size=Pt(13), color=CREAM, align=PP_ALIGN.CENTER)

# Pain points
pain = [
    "Analysts drown in PDFs — financial statements, legal contracts, market reports",
    "Coverage is inconsistent; the most junior person reads the most documents",
    "No centralised audit trail — version chaos across email and shared drives",
    "Findings are unstructured prose with no traceable citations",
    "Regulatory and legal red flags get buried in 400-page data rooms",
]
add_bullet(s, pain, Inches(0.7), Inches(3.2), Inches(11.6), Inches(3.5),
           size=Pt(16.5), color=INK, spacing=Pt(9))

section_bar(s, "The Problem")


# ── 3. OUR SOLUTION ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "Our Solution", "DealRoom AI — from documents to deal brief in hours")

# Central value proposition box
add_rect(s, Inches(0.7), Inches(1.5), Inches(11.6), Inches(1.15), fill=SAND)
add_text(s,
         "Upload your data room documents → AI agents process, research, and synthesise → "
         "Structured due diligence report with citations, risk score, and management Q&A",
         Inches(0.9), Inches(1.6), Inches(11.2), Inches(0.85),
         size=Pt(16), color=INK, align=PP_ALIGN.CENTER, font="Georgia")

# Three pillars
pillars = [
    ("Speed", "Hours, not weeks.\nThree AI agents work in parallel while analysts focus on decisions."),
    ("Coverage", "Nothing missed.\nEvery document is read, every claim is cited, every risk is flagged."),
    ("Governance", "Enterprise-grade.\nRole-based access, approval workflow, watermarking, and full audit trail."),
]
for i, (title, body) in enumerate(pillars):
    px = Inches(0.7) + i * Inches(3.9)
    add_rect(s, px, Inches(2.9), Inches(3.6), Inches(3.5), fill=WHITE,
             line=GOLD, line_w=Pt(1.5))
    # Green header bar
    add_rect(s, px, Inches(2.9), Inches(3.6), Inches(0.55), fill=GREEN)
    add_text(s, title, px + Inches(0.15), Inches(2.95), Inches(3.3), Inches(0.5),
             size=Pt(18), bold=True, color=GOLD, font="Georgia")
    add_text(s, body, px + Inches(0.2), Inches(3.55), Inches(3.3), Inches(2.7),
             size=Pt(15), color=INK)

section_bar(s, "The Solution")


# ── 4. HOW IT WORKS ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "How It Works", "Three AI agents, one governed pipeline")

# Pipeline flow
stages = [
    ("📄", "Document\nUpload", "PDF · DOCX · XLSX\nDrag-and-drop into\nyour Deal Room"),
    ("🔍", "Ingestion\nAgent", "Parse · Chunk · Embed\nStored in pgvector\nfor semantic search"),
    ("🌐", "Research\nAgent", "Web · Financial data\nNews sentiment\nSEC filings"),
    ("🧠", "Synthesis\nAgent", "Hybrid RAG retrieval\nGPT-4o report\nAll claims cited"),
    ("📊", "ML Risk\nScore", "XGBoost classifier\n4-tier risk + SHAP\nexplanations"),
    ("✅", "Structured\nReport", "6-section brief\nApproval workflow\nAudit-ready"),
]

box_w = Inches(1.82)
box_h = Inches(3.4)
gap   = Inches(0.22)
start_x = Inches(0.5)
box_y = Inches(2.0)

for i, (icon, name, detail) in enumerate(stages):
    bx = start_x + i*(box_w+gap)
    col = GREEN if i % 2 == 0 else DKGREEN
    add_rect(s, bx, box_y, box_w, box_h, fill=col)
    add_text(s, icon, bx, box_y + Inches(0.15), box_w, Inches(0.6),
             size=Pt(28), align=PP_ALIGN.CENTER)
    add_text(s, name, bx, box_y + Inches(0.75), box_w, Inches(0.65),
             size=Pt(14), bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="Georgia")
    add_text(s, detail, bx + Inches(0.1), box_y + Inches(1.45), box_w - Inches(0.2), Inches(1.9),
             size=Pt(11.5), color=CREAM, align=PP_ALIGN.CENTER)
    # Arrow
    if i < len(stages)-1:
        ax = bx + box_w + Inches(0.03)
        add_text(s, "▶", ax, box_y + Inches(1.4), gap + Inches(0.16), Inches(0.5),
                 size=Pt(14), color=GOLD, align=PP_ALIGN.CENTER)

# Note
add_text(s, "All agents are traced end-to-end via LangSmith · Redis caches research & ML inference",
         Inches(0.7), Inches(5.65), Inches(11), Inches(0.4),
         size=Pt(12), color=STONE, align=PP_ALIGN.CENTER, italic=True)

section_bar(s, "Architecture — How It Works")


# ── 5. AGENT 1: INGESTION ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "Agent 1 — Document Ingestion", "Turning raw files into searchable knowledge")

# Left column
add_text(s, "What it does:", Inches(0.7), Inches(1.55), Inches(5.5), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN)
left = [
    "Parses PDF, DOCX, and XLSX files with pdfplumber",
    "Classifies each document (financial statement, legal contract, market report…) via GPT-4o-mini",
    "Prose chunks: 512-token sliding window with 64-token overlap",
    "Tables extracted separately — each table is one atomic chunk",
    "Batch-embeds with OpenAI text-embedding-3-small (1 536 dims)",
    "Inserts into PostgreSQL with pgvector extension for semantic search",
]
add_bullet(s, left, Inches(0.7), Inches(1.95), Inches(5.8), Inches(4.5),
           size=Pt(15.5), color=INK)

# Right column — tech decisions
add_rect(s, Inches(7.0), Inches(1.4), Inches(5.5), Inches(5.2), fill=SAND)
add_text(s, "Key technical decisions", Inches(7.15), Inches(1.5), Inches(5.2), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN)
decisions = [
    ("Tables ≠ Prose", "Tables are kept as single chunks so the model receives complete rows — splitting them breaks financial context."),
    ("Async pipeline", "Runs as a FastAPI BackgroundTask on upload — user gets instant confirmation, ingestion continues silently."),
    ("Batch embedding", "100 chunks per API call — minimises OpenAI latency and cost vs. one call per chunk."),
    ("pgvector", "Vector search lives in the same Postgres instance as all other data — no separate vector DB to operate."),
]
dy = Inches(2.0)
for heading, body in decisions:
    add_text(s, f"▸ {heading}", Inches(7.2), dy, Inches(5.1), Inches(0.35),
             size=Pt(13.5), bold=True, color=DKGREEN)
    add_text(s, body, Inches(7.35), dy + Inches(0.35), Inches(4.95), Inches(0.6),
             size=Pt(12), color=INK)
    dy += Inches(1.12)

section_bar(s, "Agent 1 — Document Ingestion")


# ── 6. AGENT 2: RESEARCH ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "Agent 2 — Research Agent", "Autonomous analyst that goes beyond the data room")

add_text(s, "The problem with only reading uploaded documents:",
         Inches(0.7), Inches(1.55), Inches(11.5), Inches(0.4),
         size=Pt(14), italic=True, color=STONE)
add_text(s, "Companies don't upload bad news about themselves.",
         Inches(0.7), Inches(1.9), Inches(11.5), Inches(0.4),
         size=Pt(17), bold=True, color=GREEN, font="Georgia")

# LangGraph ReAct loop diagram
add_rect(s, Inches(0.7), Inches(2.45), Inches(5.5), Inches(3.8), fill=SAND)
add_text(s, "LangGraph ReAct Loop", Inches(0.85), Inches(2.55), Inches(5.2), Inches(0.4),
         size=Pt(13), bold=True, color=DKGREEN)

loop_steps = [
    ("reason", "LLM decides what to research next"),
    ("act", "Executes one tool call"),
    ("evaluate", "Checks if coverage is sufficient"),
]
for j, (node, desc) in enumerate(loop_steps):
    ny = Inches(3.05) + j * Inches(0.9)
    add_rect(s, Inches(0.9), ny, Inches(1.4), Inches(0.6), fill=GREEN)
    add_text(s, node, Inches(0.9), ny, Inches(1.4), Inches(0.6),
             size=Pt(13), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(2.4), ny + Inches(0.08), Inches(3.5), Inches(0.5),
             size=Pt(13), color=INK)
    if j < 2:
        add_text(s, "↓", Inches(1.55), ny + Inches(0.62), Inches(0.4), Inches(0.3),
                 size=Pt(13), color=GOLD, align=PP_ALIGN.CENTER)

add_text(s, "↺  Loop until sufficient OR 12 tool calls",
         Inches(0.9), Inches(5.82), Inches(5.0), Inches(0.4),
         size=Pt(12), color=GOLD, italic=True)

# Tools
add_text(s, "Research tools available:", Inches(6.5), Inches(2.45), Inches(6.2), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN)
tools = [
    ("🔎  Web search", "Recent news, controversies, press releases"),
    ("💰  Financial data", "Market cap, P/E, EBITDA margin, revenue growth"),
    ("📰  News sentiment", "90-day headlines with sentiment scoring"),
    ("🏢  Competitor analysis", "Market share, competitive position"),
    ("🗂️  SEC EDGAR filings", "10-K, 10-Q, 8-K regulatory filings"),
]
ty = Inches(2.95)
for icon_name, desc in tools:
    add_rect(s, Inches(6.5), ty, Inches(6.1), Inches(0.67), fill=WHITE,
             line=STONE, line_w=Pt(0.75))
    add_text(s, icon_name, Inches(6.65), ty + Inches(0.05), Inches(2.5), Inches(0.35),
             size=Pt(13), bold=True, color=DKGREEN)
    add_text(s, desc, Inches(6.65), ty + Inches(0.35), Inches(5.7), Inches(0.28),
             size=Pt(11.5), color=STONE)
    ty += Inches(0.74)

section_bar(s, "Agent 2 — Research Agent (LangGraph ReAct)")


# ── 7. AGENT 3: SYNTHESIS ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "Agent 3 — Synthesis Agent", "From knowledge to structured due diligence report")

# Left: retrieval steps
add_text(s, "Hybrid Retrieval (RAG)", Inches(0.7), Inches(1.55), Inches(5.5), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN)
retrieval = [
    ("Semantic search", "pgvector cosine similarity — finds conceptually related chunks"),
    ("BM25 keyword search", "pg_trgm + BM25 — finds exact terminology matches"),
    ("Reciprocal Rank Fusion", "Merges both ranked lists into one optimal top-15 result set"),
    ("Per-section queries", "6 sections each get their own targeted retrieval query"),
]
ry = Inches(2.0)
for heading, body in retrieval:
    add_rect(s, Inches(0.7), ry, Inches(0.45), Inches(0.45), fill=GOLD)
    add_text(s, heading, Inches(1.25), ry, Inches(4.8), Inches(0.3),
             size=Pt(13.5), bold=True, color=INK)
    add_text(s, body, Inches(1.25), ry + Inches(0.3), Inches(4.8), Inches(0.42),
             size=Pt(12), color=STONE)
    ry += Inches(0.88)

# Right: report structure
add_rect(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.3), fill=GREEN)
add_text(s, "Report Structure (GPT-4o)", Inches(6.95), Inches(1.5), Inches(5.5), Inches(0.45),
         size=Pt(14), bold=True, color=GOLD, font="Georgia")
sections = [
    "1.  Executive Summary",
    "2.  Financial Health",
    "3.  Legal Flags",
    "4.  Commercial Assessment",
    "5.  Red Flags",
    "6.  Key Questions for Management",
]
sy = Inches(2.05)
for sec in sections:
    add_text(s, sec, Inches(7.1), sy, Inches(5.3), Inches(0.5),
             size=Pt(14.5), color=CREAM)
    sy += Inches(0.55)

# Citation rule
add_rect(s, Inches(6.8), Inches(5.45), Inches(5.8), Inches(0.72), fill=DKGREEN)
add_text(s, '★  Every claim must cite its source: [filename.pdf, p.12]',
         Inches(6.95), Inches(5.52), Inches(5.5), Inches(0.6),
         size=Pt(13), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Bonus: Management Q&A
add_text(s, "Bonus: Management Q&A Generator",
         Inches(0.7), Inches(5.6), Inches(5.8), Inches(0.4),
         size=Pt(13), bold=True, color=GREEN)
add_text(s, "A single focused GPT-4o call turns red flags and financial findings into "
         "sharp, traceable management interview questions — grouped by priority (critical / high / medium).",
         Inches(0.7), Inches(5.98), Inches(5.8), Inches(0.75),
         size=Pt(12.5), color=INK)

section_bar(s, "Agent 3 — Synthesis Agent (RAG + GPT-4o)")


# ── 8. ML RISK CLASSIFIER ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "ML Financial Risk Classifier", "Quantitative signal to complement qualitative analysis")

# Top description
add_text(s,
    "Before an acquisition, analysts need a number — not just a narrative. "
    "Our XGBoost model takes 8 financial ratios and outputs a risk tier with an explanation.",
    Inches(0.7), Inches(1.5), Inches(11.5), Inches(0.6),
    size=Pt(15.5), color=INK, italic=True)

# Four risk tier boxes
tiers = [
    ("LOW", RGBColor(0x16,0x7A,0x3E), "Healthy financials\nno near-term stress"),
    ("MEDIUM", RGBColor(0xC4,0x8E,0x00), "Stable but\nworth monitoring"),
    ("HIGH", RGBColor(0xC4,0x5C,0x00), "Elevated risk —\nearly distress signals"),
    ("CRITICAL", RGBColor(0xB0,0x1F,0x1F), "High likelihood of\nfinancial failure"),
]
tx = Inches(0.7)
for label, col, desc in tiers:
    add_rect(s, tx, Inches(2.35), Inches(2.75), Inches(1.3), fill=col)
    add_text(s, label, tx, Inches(2.45), Inches(2.75), Inches(0.65),
             size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Georgia")
    add_text(s, desc, tx + Inches(0.1), Inches(3.05), Inches(2.55), Inches(0.55),
             size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)
    tx += Inches(2.88)

# Features
add_text(s, "Input: 8 financial ratios", Inches(0.7), Inches(3.85), Inches(5.8), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN)
features = [
    "Current Ratio — can it pay its bills?",
    "Debt-to-Equity — how leveraged is it?",
    "Interest Coverage — can it service its debt?",
    "EBITDA Margin — is the core business profitable?",
    "Revenue Growth YoY — growing or shrinking?",
    "Cash Burn Rate — how fast is it consuming cash?",
    "Working Capital Ratio — operational buffer",
    "Gross Margin — pricing power",
]
add_bullet(s, features, Inches(0.7), Inches(4.2), Inches(5.5), Inches(2.9),
           size=Pt(13), color=INK, spacing=Pt(3))

# SHAP explainability box
add_rect(s, Inches(6.8), Inches(3.75), Inches(5.8), Inches(2.65), fill=SAND)
add_text(s, "SHAP Explainability", Inches(6.95), Inches(3.85), Inches(5.5), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN)
add_text(s,
    'The model tells you WHY:\n\n'
    '"Scored HIGH-RISK because:\n'
    '  ▸ EBITDA margin is negative  ↑ risk\n'
    '  ▸ Debt-to-equity is very high  ↑ risk\n'
    '  ▸ Current ratio is adequate  ↓ risk"\n\n'
    'Every prediction is auditable — no black box.',
    Inches(6.95), Inches(4.3), Inches(5.4), Inches(2.0),
    size=Pt(13), color=INK)

section_bar(s, "ML Risk Classifier — XGBoost + SHAP")


# ── 9. DATA SOURCE ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "Training Data", "UCI Polish Bankruptcy Dataset")

# Big numbers
stats = [
    ("40,000+", "Annual financial statements"),
    ("7,000+", "Companies observed"),
    ("5", "Temporal slices (years)"),
    ("8", "Financial ratios per company"),
]
sx = Inches(0.7)
for stat_num, stat_label in stats:
    add_rect(s, sx, Inches(1.5), Inches(2.75), Inches(1.5), fill=GREEN)
    add_text(s, stat_num, sx, Inches(1.6), Inches(2.75), Inches(0.8),
             size=Pt(34), bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="Georgia")
    add_text(s, stat_label, sx + Inches(0.1), Inches(2.3), Inches(2.55), Inches(0.6),
             size=Pt(12), color=CREAM, align=PP_ALIGN.CENTER)
    sx += Inches(2.9)

# Source description
add_rect(s, Inches(0.7), Inches(3.2), Inches(11.6), Inches(1.0), fill=SAND)
add_text(s,
    "Source:  Emerging Markets Information Service (EMIS), 2000–2012  ·  "
    "Published via UCI Machine Learning Repository (2016)\n"
    "Format: ARFF files (64 financial ratios per company) — we extract the 8 most predictive",
    Inches(0.9), Inches(3.32), Inches(11.2), Inches(0.8),
    size=Pt(14), color=INK, italic=True)

# Label engineering
add_text(s, "Label engineering — creating 4 risk classes from 2 raw signals:",
         Inches(0.7), Inches(4.38), Inches(11.5), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN)
labels = [
    ("LOW (0)", RGBColor(0x16,0x7A,0x3E),
     "Bankrupt = No, observed 4–5 years out  →  Financially healthy, no near-term stress"),
    ("MEDIUM (1)", RGBColor(0xC4,0x8E,0x00),
     "Bankrupt = No, observed 2–3 years out  →  Stable but worth monitoring"),
    ("HIGH (2)", RGBColor(0xC4,0x5C,0x00),
     "Bankrupt = No (year 1) OR Bankrupt = Yes (year 4–5)  →  Early distress signals"),
    ("CRITICAL (3)", RGBColor(0xB0,0x1F,0x1F),
     "Bankrupt = Yes, observed 1–3 years out  →  Imminent financial failure"),
]
ly = Inches(4.88)
for lbl, col, desc in labels:
    add_rect(s, Inches(0.7), ly, Inches(1.65), Inches(0.42), fill=col)
    add_text(s, lbl, Inches(0.72), ly + Inches(0.02), Inches(1.6), Inches(0.4),
             size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(2.5), ly + Inches(0.04), Inches(9.8), Inches(0.38),
             size=Pt(13), color=INK)
    ly += Inches(0.5)

# Performance
add_text(s,
    "Achieved Macro F1: 0.55  ·  Model catches ~46% of genuinely critical companies "
    "(treated as a first-pass filter, not a final verdict)  ·  "
    "Top SHAP drivers: EBITDA margin · Interest coverage · Revenue growth",
    Inches(0.7), Inches(6.9), Inches(11.5), Inches(0.45),
    size=Pt(11.5), color=STONE, italic=True)

section_bar(s, "Training Data — UCI Polish Bankruptcy Dataset · EMIS 2000–2012")


# ── 10. THE PLATFORM ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "The Platform", "Enterprise-grade SaaS for consulting firms")

# Left: feature list
add_text(s, "Core Platform Features", Inches(0.7), Inches(1.55), Inches(5.8), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN)
platform_features = [
    "Multi-tenant — every consulting firm gets a completely isolated workspace",
    "Deal Rooms — dedicated workspace per acquisition target with explicit member invitations",
    "Role-based access — Admin vs. Analyst with different permissions per deal",
    "Document permissions — per-document access control list on top of deal room membership",
    "Approval workflow — Draft → In Review → Approved (signed by senior analyst)",
    "Annotation system — inline comments, dispute flags, and resolution threads",
    "Management Q&A — auto-generated interview questions from red flags",
    "Watermarked exports — every downloaded document is stamped with analyst identity",
    "Append-only audit trail — every action logged with timestamp, actor, and payload",
    "Email notifications — async SMTP dispatch on report events",
]
add_bullet(s, platform_features, Inches(0.7), Inches(1.98), Inches(5.8), Inches(4.7),
           size=Pt(13.5), color=INK, spacing=Pt(4))

# Right: workflow diagram
add_rect(s, Inches(7.0), Inches(1.4), Inches(5.6), Inches(5.4), fill=SAND)
add_text(s, "Report Lifecycle", Inches(7.15), Inches(1.5), Inches(5.3), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN)
workflow = [
    (RGBColor(0x88,0x88,0x88), "DRAFT", "AI generates report\nanalyst reviews & edits"),
    (RGBColor(0x1A,0x5E,0x3A), "IN REVIEW", "Submitted by analyst\nsenior reviews findings"),
    (RGBColor(0xD4,0xA8,0x4B), "APPROVED", "Signed off by senior\nExport unlocked"),
]
wy = Inches(2.05)
for col, status, desc in workflow:
    add_rect(s, Inches(7.2), wy, Inches(2.0), Inches(1.2), fill=col)
    add_text(s, status, Inches(7.2), wy + Inches(0.2), Inches(2.0), Inches(0.5),
             size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, Inches(9.35), wy + Inches(0.15), Inches(3.0), Inches(0.9),
             size=Pt(13), color=INK)
    if col != RGBColor(0xD4,0xA8,0x4B):
        add_text(s, "↓", Inches(8.1), wy + Inches(1.22), Inches(0.5), Inches(0.35),
                 size=Pt(16), color=GREEN, align=PP_ALIGN.CENTER)
    wy += Inches(1.58)

# Storage note
add_rect(s, Inches(7.0), Inches(5.48), Inches(5.6), Inches(0.65), fill=DKGREEN)
add_text(s, "Documents stored in MinIO (S3-compatible)\nPostgreSQL + pgvector for all structured data",
         Inches(7.1), Inches(5.54), Inches(5.3), Inches(0.58),
         size=Pt(12), color=CREAM, align=PP_ALIGN.CENTER)

section_bar(s, "Platform — Multi-Tenant SaaS with Governed Workflows")


# ── 11. TECH STACK ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "Technology Stack", "Production-grade, cloud-ready, open-source core")

categories = [
    ("Frontend",   ["React 18 + Vite", "Tailwind CSS", "React Router v6", "Zustand (auth state)"]),
    ("Backend",    ["FastAPI (Python 3.12)", "SQLAlchemy (async)", "Alembic migrations", "JWT + bcrypt auth"]),
    ("AI / LLM",  ["LangGraph (agent graphs)", "LangChain tools", "GPT-4o (synthesis)", "GPT-4o-mini (classify)"]),
    ("ML",         ["scikit-learn pipeline", "XGBoost classifier", "SHAP explainability", "Macro F1: 0.55"]),
    ("Data",       ["PostgreSQL 16 + pgvector", "OpenAI embeddings", "BM25 + RRF retrieval", "pdfplumber parsing"]),
    ("Infra",      ["Docker + Docker Compose", "Redis + ARQ workers", "MinIO object storage", "GitHub Actions CI"]),
]

cols = 3
rows = 2
box_w = Inches(3.9)
box_h = Inches(2.35)
bx0   = Inches(0.55)
by0   = Inches(1.45)
col_gap = Inches(0.2)
row_gap = Inches(0.18)

for i, (cat, items) in enumerate(categories):
    col_i = i % cols
    row_i = i // cols
    bx = bx0 + col_i * (box_w + col_gap)
    by = by0 + row_i * (box_h + row_gap)
    add_rect(s, bx, by, box_w, box_h, fill=WHITE, line=STONE, line_w=Pt(0.8))
    add_rect(s, bx, by, box_w, Inches(0.46), fill=GREEN)
    add_text(s, cat, bx + Inches(0.15), by + Inches(0.06), box_w - Inches(0.2), Inches(0.38),
             size=Pt(14), bold=True, color=GOLD, font="Georgia")
    for j, item in enumerate(items):
        add_text(s, f"▸  {item}",
                 bx + Inches(0.18), by + Inches(0.54) + j*Inches(0.42),
                 box_w - Inches(0.25), Inches(0.4),
                 size=Pt(13), color=INK)

section_bar(s, "Technology Stack")


# ── 12. KEY DIFFERENTIATORS ───────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)

slide_title(s, "Why DealRoom AI?", "Built for how M&A due diligence actually works")

diffs = [
    ("⚡  Speed", "Hours to a full brief", "vs. 6–12 weeks of manual analysis.\nFirms can review more deals with the same team."),
    ("📌  Traceability", "Every claim is cited", "No AI hallucination goes undetected — every sentence links back to its source page and chunk."),
    ("🔒  Security", "Enterprise access control", "Per-deal membership, per-document ACLs, watermarked exports, and an append-only audit trail."),
    ("🤖  Explainable ML", "Not a black box", "SHAP values explain exactly which financial ratios drove the risk score — auditable by regulators."),
    ("🏗️  Production-ready", "Built to scale", "Async worker queue, Redis caching, Docker Compose deployment, CI/CD with automated ML gating."),
    ("🎯  Analyst-in-the-loop", "AI assists, humans decide", "Annotation threads, dispute flags, and the approval workflow ensure a human signs off every report."),
]

dx = Inches(0.55)
dy_base = Inches(1.55)
d_w = Inches(3.85)
d_h = Inches(2.4)

for i, (icon_label, headline, body) in enumerate(diffs):
    col_i = i % 3
    row_i = i // 3
    bx = dx + col_i * (d_w + Inches(0.2))
    by = dy_base + row_i * (d_h + Inches(0.18))
    add_rect(s, bx, by, d_w, d_h, fill=SAND)
    add_rect(s, bx, by, d_w, Inches(0.48), fill=DKGREEN)
    add_text(s, icon_label, bx + Inches(0.12), by + Inches(0.07), d_w, Inches(0.38),
             size=Pt(13.5), bold=True, color=GOLD)
    add_text(s, headline, bx + Inches(0.12), by + Inches(0.6), d_w - Inches(0.2), Inches(0.42),
             size=Pt(15), bold=True, color=GREEN, font="Georgia")
    add_text(s, body, bx + Inches(0.12), by + Inches(1.0), d_w - Inches(0.2), Inches(1.3),
             size=Pt(13), color=INK)

section_bar(s, "Differentiators")


# ── 13. LIVE DEMO ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, GREEN)

# Full-bleed green slide
add_rect(s, Inches(0), Inches(0), W, H, fill=GREEN)
add_rect(s, Inches(0), H - Inches(0.12), W, Inches(0.12), fill=GOLD)

# Large "DEMO" text
add_text(s, "LIVE DEMO", Inches(1), Inches(1.5), Inches(11.5), Inches(2.5),
         size=Pt(96), bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="Georgia")

gold_line(s, Inches(3), Inches(4.0), Inches(7))

demo_steps = [
    "1.  Create a Deal Room and invite members",
    "2.  Upload documents — watch ingestion in real time",
    "3.  Trigger analysis — all three agents run",
    "4.  Review the structured report with citations",
    "5.  See the ML risk score and SHAP factors",
    "6.  Add annotations, submit for approval",
]
add_bullet(s, demo_steps, Inches(2.2), Inches(4.3), Inches(9), Inches(2.8),
           size=Pt(17), color=CREAM, bullet_color=GOLD)


# ── 14. Q&A / THANK YOU ───────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, CREAM)
add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=GREEN)
add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=GOLD)

# Logo icon small
draw_logo(s, Inches(6.67), Inches(2.2), scale=1.3)

add_text(s, "DealRoom", Inches(7.6), Inches(1.4), Inches(4.8), Inches(0.9),
         size=Pt(44), bold=True, color=INK, font="Georgia")
add_text(s, "AI", Inches(10.88), Inches(1.4), Inches(2), Inches(0.9),
         size=Pt(44), bold=False, color=GREEN, font="Georgia")

gold_line(s, Inches(7.6), Inches(2.35), Inches(4.0))
add_text(s, "DUE DILIGENCE.  REIMAGINED.", Inches(7.6), Inches(2.52), Inches(5.5), Inches(0.4),
         size=Pt(11), bold=True, color=INK)

add_text(s, "Thank You", Inches(0.7), Inches(2.0), Inches(5.5), Inches(1.4),
         size=Pt(64), bold=True, color=GREEN, font="Georgia")
add_text(s, "Questions & Discussion", Inches(0.7), Inches(3.4), Inches(5.5), Inches(0.6),
         size=Pt(22), color=STONE, font="Georgia")

gold_line(s, Inches(0.7), Inches(4.1), Inches(4.5))

# Contact / links
add_text(s, "oceansmelody.om@gmail.com",
         Inches(0.7), Inches(4.3), Inches(5.5), Inches(0.45),
         size=Pt(15), color=INK)
add_text(s, "github.com/sShawraba  ·  DealRoom AI",
         Inches(0.7), Inches(4.75), Inches(5.5), Inches(0.45),
         size=Pt(14), color=STONE)

# Quick recap
add_rect(s, Inches(0.7), Inches(5.5), Inches(11.6), Inches(1.0), fill=SAND)
recap = (
    "3 AI agents  ·  XGBoost risk classifier trained on 40k+ records  ·  "
    "Hybrid RAG retrieval  ·  Full citation trail  ·  Multi-tenant SaaS  ·  "
    "Governed approval workflow"
)
add_text(s, recap, Inches(0.9), Inches(5.65), Inches(11.2), Inches(0.75),
         size=Pt(13.5), color=INK, align=PP_ALIGN.CENTER, italic=True)


# ── SAVE ───────────────────────────────────────────────────────────────────
output = "/home/soup/due-diligence-ai/DealRoom_AI_Presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
